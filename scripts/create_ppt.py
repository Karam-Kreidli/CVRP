"""
Generate a PowerPoint presentation explaining the CVRP RL solver project.
Run: python scripts/create_ppt.py
Output: CVRP_RL_Solver.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── colour palette ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x23, 0x3A)   # slide backgrounds / headers
MID_BLUE    = RGBColor(0x1F, 0x4E, 0x79)   # accent bars
ACCENT_CYAN = RGBColor(0x00, 0xB0, 0xD4)   # highlights
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)   # text boxes
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
ORANGE      = RGBColor(0xFF, 0x85, 0x00)

# ── helpers ─────────────────────────────────────────────────────────────────

def inches(*args):
    return [Inches(a) for a in args]


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_bullet_box(slide, bullets, left, top, width, height,
                   title=None, title_size=16, bullet_size=13,
                   bg_color=None, text_color=WHITE):
    if bg_color:
        add_rect(slide, left, top, width, height, bg_color)
    txb = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1),
                                   Inches(width - 0.3), Inches(height - 0.2))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size = Pt(title_size)
        r.font.bold = True
        r.font.color.rgb = ACCENT_CYAN
        tf.add_paragraph()  # spacer

    first = True
    for b in bullets:
        if title or not first:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]
            first = False
        para.alignment = PP_ALIGN.LEFT
        para.space_before = Pt(2)
        r = para.add_run()
        r.text = b
        r.font.size = Pt(bullet_size)
        r.font.color.rgb = text_color


# ── slide builders ──────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 – title / cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, DARK_BLUE)

    # accent bar top
    add_rect(slide, 0, 0, 10, 0.08, ACCENT_CYAN)
    # accent bar bottom
    add_rect(slide, 0, 7.42, 10, 0.08, ACCENT_CYAN)

    # main title
    add_textbox(slide, "RL-Guided CVRP Solver",
                0.5, 1.2, 9, 1.2,
                font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # subtitle
    add_textbox(slide, "Capacitated Vehicle Routing Problem  ·  GECCO 2026 ML4VRP Competition",
                0.5, 2.5, 9, 0.7,
                font_size=18, bold=False, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

    # divider
    add_rect(slide, 3.5, 3.35, 3, 0.04, ACCENT_CYAN)

    # tagline
    add_textbox(slide,
                "A Reinforcement Learning agent learns when and how to push\n"
                "a state-of-the-art CVRP solver for minimal fleet & distance.",
                0.7, 3.6, 8.6, 1.4,
                font_size=16, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # pills
    for i, (label, x) in enumerate([
        ("PPO · Actor-Critic", 0.8),
        ("HGS-CVRP (hygese)", 3.7),
        ("X-Dataset · 59 instances", 6.5),
    ]):
        add_rect(slide, x, 5.3, 2.5, 0.45, MID_BLUE)
        add_textbox(slide, label, x + 0.05, 5.32, 2.4, 0.4,
                    font_size=13, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    """Slide 2 – problem statement."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, 10, 0.08, ACCENT_CYAN)

    add_textbox(slide, "The Problem: CVRP", 0.4, 0.18, 9, 0.65,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.88, 1.5, 0.04, ACCENT_CYAN)

    # two columns
    add_bullet_box(slide,
        ["Route a fleet of vehicles from a central depot",
         "Each customer has a location and demand",
         "Each vehicle has a fixed capacity",
         "All customers must be served exactly once",
         "Minimise: cost = 1000×NV + Total Distance"],
        left=0.4, top=1.1, width=4.5, height=3.5,
        title="What is CVRP?", bg_color=MID_BLUE)

    add_bullet_box(slide,
        ["NV (vehicles) is multiplied by 1000",
         "→ removing one truck saves 1000 points",
         "Distance only matters after fleet is minimised",
         "N-P Hard: exact methods fail at scale",
         "HGS-CVRP (C++) is state-of-the-art heuristic"],
        left=5.2, top=1.1, width=4.4, height=3.5,
        title="Competition Objective", bg_color=MID_BLUE)

    # formula box
    add_rect(slide, 1.5, 4.9, 7, 0.75, RGBColor(0x0D, 0x2B, 0x55))
    add_textbox(slide, "Score  =  1000 × NV  +  Total Distance  →  minimise",
                1.6, 4.95, 6.8, 0.65,
                font_size=20, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)


def slide_approach(prs):
    """Slide 3 – high-level approach."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, 10, 0.08, ACCENT_CYAN)

    add_textbox(slide, "Our Approach: RL-Guided Search", 0.4, 0.18, 9, 0.65,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.88, 2.0, 0.04, ACCENT_CYAN)

    add_textbox(slide,
        "Instead of solving CVRP directly with neural networks, we train a lightweight "
        "RL agent to make high-level strategic decisions—then let the world-class HGS "
        "solver do the heavy lifting.",
        0.4, 1.05, 9.2, 0.9, font_size=15, color=LIGHT_GREY)

    # pipeline boxes
    boxes = [
        ("1  Instance\nFeatures", "7 hand-crafted\nspatial stats", 0.3),
        ("2  RL Agent\n(Fleet Manager)", "Actor-Critic\n5 700 params", 2.55),
        ("3  HGS-CVRP\nSolver", "C++ high-perf\nheuristic", 4.8),
        ("4  Reward\nSignal", "% improvement\nover episode best", 7.05),
    ]
    for title, sub, x in boxes:
        add_rect(slide, x, 2.15, 2.1, 1.5, MID_BLUE)
        add_textbox(slide, title, x + 0.08, 2.2, 1.95, 0.75,
                    font_size=13, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
        add_textbox(slide, sub, x + 0.08, 2.95, 1.95, 0.65,
                    font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # arrows between boxes
    for ax in [2.42, 4.67, 6.92]:
        add_textbox(slide, "→", ax, 2.65, 0.3, 0.5,
                    font_size=20, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

    # episode description
    add_bullet_box(slide,
        ["Episode = one CVRP instance × 50 decision steps",
         "Each step: agent picks an action → HGS runs → new best tracked",
         "PPO updates policy using cumulative reward",
         "Curriculum: 20 epochs on small (N≤100) then all 59 instances"],
        left=0.4, top=3.95, width=9.2, height=2.4,
        title="Training Loop", bg_color=MID_BLUE)


def slide_agent(prs):
    """Slide 4 – the RL agent & actions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, 10, 0.08, ACCENT_CYAN)

    add_textbox(slide, "Fleet Manager – RL Agent", 0.4, 0.18, 9, 0.65,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.88, 1.8, 0.04, ACCENT_CYAN)

    # architecture box
    add_bullet_box(slide,
        ["Input: 14-dim vector  (7 instance features + 7 live solver stats)",
         "Hidden: two 64-unit ReLU layers  (shared trunk)",
         "Actor head: 64 → 7 action logits  (+ action masking)",
         "Critic head: 64 → 1 state-value estimate",
         "~5,700 parameters total  –  tiny by design"],
        left=0.4, top=1.05, width=4.5, height=3.2,
        title="Network Architecture", bg_color=MID_BLUE)

    # actions table
    add_rect(slide, 5.1, 1.05, 4.6, 3.2, MID_BLUE)
    add_textbox(slide, "7 Discrete Actions", 5.2, 1.08, 4.4, 0.4,
                font_size=14, bold=True, color=ACCENT_CYAN)
    actions = [
        ("#0  FREE_SAME",   "Unconstrained, same seed"),
        ("#1  FREE_NEW",    "Unconstrained, new seed"),
        ("#2  LOCK_SAME",   "Lock best NV, same seed"),
        ("#3  LOCK_NEW",    "Lock best NV, new seed"),
        ("#4  PUSH_SAME",   "Try NV−1, same seed  ×2 iters"),
        ("#5  PUSH_NEW",    "Try NV−1, new seed  ×2 iters"),
        ("#6  FORCE_MIN",   "Force NV_min  ×3 iters"),
    ]
    for i, (a, d) in enumerate(actions):
        y = 1.58 + i * 0.35
        add_textbox(slide, a, 5.2, y, 2.1, 0.33,
                    font_size=10, bold=True, color=ACCENT_CYAN)
        add_textbox(slide, d, 7.3, y, 2.3, 0.33,
                    font_size=10, color=LIGHT_GREY)

    # observation box
    add_bullet_box(slide,
        ["Instance features (computed once):  size, demand fill ratio, "
         "mean/std distance, depot centrality, demand CV, capacity tightness",
         "Live solver stats (updated each step):  best NV, best TD, best score, "
         "step fraction, NV fraction, score improvement, NV headroom"],
        left=0.4, top=4.45, width=9.2, height=2.1,
        title="Observation Space (14 dims)", bg_color=MID_BLUE)

    add_textbox(slide, "* Actions 4-6 masked when NV ≤ NV_min (impossible fleet reductions)",
                0.4, 6.7, 9.2, 0.35, font_size=11, color=RGBColor(0xAA, 0xAA, 0xAA))


def slide_features(prs):
    """Slide 5 – instance features."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, 10, 0.08, ACCENT_CYAN)

    add_textbox(slide, "Instance Feature Extraction", 0.4, 0.18, 9, 0.65,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.88, 2.0, 0.04, ACCENT_CYAN)

    add_textbox(slide,
        "Seven deterministic, hand-crafted features capture the spatial and structural "
        "properties of each instance without requiring a learned encoder.",
        0.4, 1.05, 9.2, 0.7, font_size=14, color=LIGHT_GREY)

    features = [
        ("size_norm",          "num_customers / 400",                    "Instance scale"),
        ("demand_fill_ratio",  "total_demand / (NV_min × capacity)",     "Vehicle packing tightness"),
        ("mean_dist_norm",     "mean_distance / max_distance",           "Average inter-customer spacing"),
        ("std_dist_norm",      "std_distance / max_distance",            "Distance variance (clustered vs spread)"),
        ("depot_centrality",   "mean_depot_dist / max_distance",         "How central the depot is"),
        ("demand_cv",          "std_demand / mean_demand",               "Demand heterogeneity"),
        ("capacity_tightness", "max_demand / capacity",                  "Hardest single-customer fill"),
    ]

    header_y = 1.85
    add_rect(slide, 0.35, header_y, 9.3, 0.38, MID_BLUE)
    for label, col_x, col_w in [("Feature", 0.45, 2.1), ("Formula", 2.75, 3.6), ("Meaning", 6.55, 3.0)]:
        add_textbox(slide, label, col_x, header_y + 0.04, col_w, 0.3,
                    font_size=12, bold=True, color=ACCENT_CYAN)

    for i, (name, formula, meaning) in enumerate(features):
        row_y = 2.28 + i * 0.55
        bg = RGBColor(0x14, 0x2A, 0x4A) if i % 2 == 0 else RGBColor(0x1A, 0x34, 0x58)
        add_rect(slide, 0.35, row_y, 9.3, 0.52, bg)
        add_textbox(slide, name,    0.45, row_y + 0.08, 2.15, 0.38, font_size=11, bold=True,  color=ACCENT_CYAN)
        add_textbox(slide, formula, 2.75, row_y + 0.08, 3.6,  0.38, font_size=11,             color=LIGHT_GREY)
        add_textbox(slide, meaning, 6.55, row_y + 0.08, 3.0,  0.38, font_size=11,             color=WHITE)


def slide_training(prs):
    """Slide 6 – PPO training."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, 10, 0.08, ACCENT_CYAN)

    add_textbox(slide, "PPO Training Loop", 0.4, 0.18, 9, 0.65,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.88, 1.4, 0.04, ACCENT_CYAN)

    add_bullet_box(slide,
        ["Algorithm: Proximal Policy Optimisation (PPO) with clipped surrogate",
         "Rollout: 8 episodes × 50 steps = 400 transitions per epoch",
         "GAE-λ advantages: λ = 0.90  (bias-variance balance)",
         "Clip ratio: ε = 0.20  (prevents large policy jumps)",
         "Entropy bonus: 0.02  (encourages continued exploration)",
         "Learning rate: 1e-4 → 5e-5  (linear decay)",
         "FP16 mixed precision on GPU  (halves memory usage)"],
        left=0.4, top=1.05, width=4.5, height=4.1,
        title="Hyper-parameters", bg_color=MID_BLUE)

    add_bullet_box(slide,
        ["Epochs 1-20: small instances only (N ≤ 100 customers)",
         "  → agent learns basic strategies quickly",
         "Epochs 21+: all 59 X-dataset instances (up to 400 customers)",
         "  → policy generalises to harder problems",
         "",
         "Evaluation every epoch on 5 fixed instances:",
         "  X-n101, X-n157, X-n223, X-n261, X-n313",
         "  (greedy policy, no exploration noise)"],
        left=5.1, top=1.05, width=4.5, height=4.1,
        title="Curriculum & Evaluation", bg_color=MID_BLUE)

    add_bullet_box(slide,
        ["Reward = % improvement over episode-best score  (positive only on new bests)",
         "Fleet-explosion penalty: −5.0  |  No-improvement: −0.5"],
        left=0.4, top=5.35, width=9.2, height=1.15,
        title="Reward Signal", bg_color=RGBColor(0x0D, 0x2B, 0x55))


def slide_architecture(prs):
    """Slide 7 – system architecture diagram (text-based)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, 10, 0.08, ACCENT_CYAN)

    add_textbox(slide, "System Architecture", 0.4, 0.18, 9, 0.65,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.88, 1.5, 0.04, ACCENT_CYAN)

    # boxes with connecting arrows (left to right, then loop back)
    components = [
        (".vrp File\n(X-dataset)", 0.3,  1.2, 1.7, 1.4),
        ("Feature\nExtractor\n(7 stats)",  2.35, 1.2, 1.7, 1.4),
        ("Fleet Manager\n(Actor-Critic\nNN)",  4.4,  1.2, 1.7, 1.4),
        ("HGS-CVRP\nSolver\n(hygese C++)",  6.45, 1.2, 1.7, 1.4),
        ("Best Solution\n& Reward",          8.1,  1.2, 1.6, 1.4),
    ]
    for label, lx, ly, lw, lh in components:
        add_rect(slide, lx, ly, lw, lh, MID_BLUE)
        add_textbox(slide, label, lx + 0.05, ly + 0.2, lw - 0.1, lh - 0.3,
                    font_size=12, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

    # arrows
    for ax in [2.07, 4.12, 6.17]:
        add_textbox(slide, "→", ax, 1.62, 0.25, 0.5,
                    font_size=18, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

    # feedback arrow label
    add_rect(slide, 1.0, 3.0, 7.8, 0.5, RGBColor(0x0D, 0x2B, 0x55))
    add_textbox(slide, "PPO update  ←  GAE advantages  ←  reward signal  (loops each epoch)",
                1.1, 3.03, 7.6, 0.4,
                font_size=12, color=ORANGE, align=PP_ALIGN.CENTER)

    # file layout
    add_bullet_box(slide,
        ["src/agent_manager.py   –  Fleet Manager (Actor-Critic NN)",
         "src/solver_engine.py   –  CVRPEnv (Gymnasium + HGS wrapper)",
         "src/train.py           –  PPO loop, GAE, checkpointing",
         "src/main.py            –  CLI entry point & smoke tests",
         "scripts/infer.py       –  Run trained agent on new instances",
         "data/*.vrp             –  59 X-dataset instances",
         "logs/best_model.pth    –  Trained checkpoint (~80 KB)"],
        left=0.4, top=3.7, width=9.2, height=3.1,
        title="File Layout", bg_color=MID_BLUE)


def slide_dataset(prs):
    """Slide 8 – dataset & results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, 10, 0.08, ACCENT_CYAN)

    add_textbox(slide, "Dataset & Benchmarks", 0.4, 0.18, 9, 0.65,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.88, 1.6, 0.04, ACCENT_CYAN)

    add_bullet_box(slide,
        ["X-dataset (Uchoa et al., 2014)  –  standard CVRP benchmark",
         "59 instances:  X-n101-k25  …  X-n401-k29",
         "100–400 customers per instance",
         "TSPLIB format: coordinates, demands, vehicle capacity",
         "Best Known Solutions (BKS) available for each instance"],
        left=0.4, top=1.05, width=4.5, height=2.8,
        title="Dataset", bg_color=MID_BLUE)

    # BKS table
    add_rect(slide, 5.1, 1.05, 4.6, 2.8, MID_BLUE)
    add_textbox(slide, "Typical BKS Score Ranges", 5.2, 1.08, 4.4, 0.4,
                font_size=14, bold=True, color=ACCENT_CYAN)
    bks = [
        ("Small  (100-150 cust.)", "≈ 20 000 – 90 000"),
        ("Medium (150-250 cust.)", "≈ 35 000 – 190 000"),
        ("Large  (250-400 cust.)", "≈ 39 000 – 241 000"),
    ]
    for i, (sz, sc) in enumerate(bks):
        y = 1.6 + i * 0.65
        add_rect(slide, 5.2, y, 4.4, 0.55,
                 RGBColor(0x14, 0x2A, 0x4A) if i % 2 == 0 else RGBColor(0x1A, 0x34, 0x58))
        add_textbox(slide, sz, 5.3, y + 0.1, 2.5, 0.38, font_size=11, color=WHITE)
        add_textbox(slide, sc, 7.9, y + 0.1, 1.6, 0.38, font_size=11, bold=True, color=ACCENT_CYAN)

    # eval instances
    add_bullet_box(slide,
        ["Fixed validation set (5 instances, evaluated every epoch):",
         "  • X-n101-k25    (100 customers)",
         "  • X-n157-k13    (156 customers)",
         "  • X-n223-k34    (222 customers)",
         "  • X-n261-k13    (260 customers)",
         "  • X-n313-k71    (312 customers)"],
        left=0.4, top=4.1, width=9.2, height=2.4,
        title="Evaluation Instances", bg_color=RGBColor(0x0D, 0x2B, 0x55))

    add_textbox(slide,
        "Training log saved to  logs/training_metrics.csv  (epoch, eval_score, NV, TD, entropy, …)",
        0.4, 6.7, 9.2, 0.35, font_size=11, color=RGBColor(0xAA, 0xAA, 0xAA))


def slide_key_insights(prs):
    """Slide 9 – key design insights."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, 10, 0.08, ACCENT_CYAN)

    add_textbox(slide, "Key Design Insights", 0.4, 0.18, 9, 0.65,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.88, 1.5, 0.04, ACCENT_CYAN)

    insights = [
        ("Hand-crafted features  (not GNN)",
         "Instant signal without training overhead. Interpretable and cheap to compute once per episode."),
        ("Fleet-target action space",
         "Directly controls HGS parameters. Three dimensions: fleet target × seed × iteration budget."),
        ("No warm-starting  (fresh solves)",
         "Each step is an independent HGS run. Avoids local optima and simplifies credit assignment."),
        ("Percentage-based reward",
         "Normalises across instance sizes. Positive only when a new best is found—unambiguous signal."),
        ("Curriculum learning",
         "Start small for rapid skill acquisition; then generalise. Prevents early over-fitting to easy instances."),
        ("FP16 mixed precision",
         "Halves GPU memory. Enables larger effective batch sizes with negligible accuracy cost."),
    ]

    for i, (title, detail) in enumerate(insights):
        col = i % 2
        row = i // 2
        lx = 0.35 + col * 4.9
        ly = 1.1 + row * 2.0
        add_rect(slide, lx, ly, 4.6, 1.75, MID_BLUE)
        add_textbox(slide, title,  lx + 0.12, ly + 0.1,  4.35, 0.45,
                    font_size=13, bold=True, color=ACCENT_CYAN)
        add_textbox(slide, detail, lx + 0.12, ly + 0.58, 4.35, 1.1,
                    font_size=11, color=LIGHT_GREY)


def slide_usage(prs):
    """Slide 10 – how to run."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, 10, 0.08, ACCENT_CYAN)

    add_textbox(slide, "How to Run", 0.4, 0.18, 9, 0.65,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.88, 0.9, 0.04, ACCENT_CYAN)

    # code blocks (dark bg)
    code_bg = RGBColor(0x0A, 0x14, 0x28)

    def code_box(slide, lines, lx, ly, lw, lh, label):
        add_rect(slide, lx, ly, lw, 0.35, RGBColor(0x0D, 0x2B, 0x55))
        add_textbox(slide, label, lx + 0.1, ly + 0.04, lw - 0.2, 0.28,
                    font_size=12, bold=True, color=ACCENT_CYAN)
        add_rect(slide, lx, ly + 0.35, lw, lh, code_bg)
        txb = slide.shapes.add_textbox(
            Inches(lx + 0.15), Inches(ly + 0.42),
            Inches(lw - 0.3), Inches(lh - 0.15))
        txb.word_wrap = False
        tf = txb.text_frame
        tf.word_wrap = False
        first = True
        for line in lines:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = line
            r.font.size = Pt(11)
            r.font.color.rgb = RGBColor(0x90, 0xEE, 0x90)  # light green
            r.font.name = "Courier New"

    code_box(slide,
        ["pip install -r requirements.txt",
         "# needs: torch  numpy  gymnasium  hygese",
         "# hygese requires a C++ compiler"],
        0.35, 1.05, 4.5, 0.95, "Installation")

    code_box(slide,
        ["# Run smoke tests (verify everything works)",
         "python -m src.main"],
        5.15, 1.05, 4.5, 0.95, "Smoke Tests")

    code_box(slide,
        ["python -m src.main train \\",
         "  --instance_path data/ \\",
         "  --epochs 50 \\",
         "  --episodes_per_epoch 8 \\",
         "  --batch_size 128 \\",
         "  --fp16 \\",
         "  --curriculum_epochs 20"],
        0.35, 2.2, 4.5, 2.2, "Training")

    code_box(slide,
        ["# Single instance",
         "python -m scripts.infer \\",
         "  --instance data/X-n200-k36.vrp",
         "",
         "# All instances (with baseline)",
         "python -m scripts.infer \\",
         "  --instance_dir data/ \\",
         "  --baseline --verbose"],
        5.15, 2.2, 4.5, 2.2, "Inference")

    add_bullet_box(slide,
        ["--resume            Resume from checkpoint",
         "--start_epoch N     Continue from epoch N",
         "--fp16              Enable mixed precision (CUDA required)",
         "--curriculum_epochs N   Epochs restricted to small instances"],
        left=0.35, top=4.6, width=9.3, height=1.85,
        title="Key CLI Arguments", bg_color=MID_BLUE)


def slide_summary(prs):
    """Slide 11 – summary / closing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, 10, 0.08, ACCENT_CYAN)
    add_rect(slide, 0, 7.42, 10, 0.08, ACCENT_CYAN)

    add_textbox(slide, "Summary", 0.4, 0.18, 9, 0.65,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.88, 0.8, 0.04, ACCENT_CYAN)

    # 4 take-away cards
    cards = [
        ("Problem",
         "CVRP: route vehicles minimising\n1000×NV + Total Distance"),
        ("Method",
         "Lightweight PPO agent (5.7k params)\nguides the powerful HGS-CVRP solver"),
        ("Key Idea",
         "Agent learns when to push fleet size\nvs. when to optimise distance"),
        ("Result",
         "Curriculum + hand-crafted features\nenable fast, generalising training"),
    ]
    for i, (title, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        lx = 0.4 + col * 4.85
        ly = 1.1 + row * 2.1
        add_rect(slide, lx, ly, 4.55, 1.85, MID_BLUE)
        add_rect(slide, lx, ly, 4.55, 0.4, ACCENT_CYAN)
        add_textbox(slide, title, lx + 0.1, ly + 0.05, 4.35, 0.32,
                    font_size=15, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        add_textbox(slide, body, lx + 0.1, ly + 0.5, 4.35, 1.25,
                    font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        "GECCO 2026 ML4VRP Competition  ·  X-dataset (Uchoa et al., 2014)  ·  HGS-CVRP (Vidal, 2022)",
        0.4, 6.95, 9.2, 0.4,
        font_size=11, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)


# ── main ────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_approach(prs)
    slide_agent(prs)
    slide_features(prs)
    slide_training(prs)
    slide_architecture(prs)
    slide_dataset(prs)
    slide_key_insights(prs)
    slide_usage(prs)
    slide_summary(prs)

    out = "CVRP_RL_Solver.pptx"
    prs.save(out)
    print(f"Saved → {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
